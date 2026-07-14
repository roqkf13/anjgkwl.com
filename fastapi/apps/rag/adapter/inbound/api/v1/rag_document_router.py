import logging
import re
from io import BytesIO

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile

from rag.adapter.inbound.api.schemas.rag_document_schema import RagUploadResponseSchema
from rag.app.ports.input.rag_document_use_case import RagDocumentUseCase
from rag.dependencies.rag_document_provider import get_rag_document_use_case

logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = (
    ".pdf",
    ".txt",
    ".md",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
    ".html",
    ".htm",
    ".json",
    ".sql",
)

rag_document_router = APIRouter(tags=["rag"])


def _extract_text(filename: str, content: bytes) -> str:
    lower = filename.lower()

    if lower.endswith(".pdf"):
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if lower.endswith(".docx"):
        from docx import Document

        doc = Document(BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs)

    if lower.endswith(".pptx"):
        from pptx import Presentation

        prs = Presentation(BytesIO(content))
        lines = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        return "\n".join(lines)

    if lower.endswith(".xlsx"):
        from openpyxl import load_workbook

        wb = load_workbook(BytesIO(content), data_only=True, read_only=True)
        lines = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append(", ".join(cells))
        return "\n".join(lines)

    if lower.endswith((".html", ".htm")):
        raw = content.decode("utf-8", errors="ignore")
        stripped = re.sub(r"<[^>]+>", " ", raw)
        return re.sub(r"\s+", " ", stripped).strip()

    return content.decode("utf-8", errors="ignore")


@rag_document_router.post(
    "/upload",
    response_model=RagUploadResponseSchema,
    summary="RAG 문서(pdf, txt, md, docx, pptx, xlsx, csv, html, json, sql) 업로드",
)
async def upload_document(
    file: UploadFile = File(...),
    use_case: RagDocumentUseCase = Depends(get_rag_document_use_case),
) -> RagUploadResponseSchema:
    filename = file.filename or ""
    if not filename.lower().endswith(ALLOWED_EXTENSIONS):
        raise HTTPException(
            status_code=400,
            detail=(
                "pdf, txt, md, docx, pptx, xlsx, csv, html, json, sql "
                "파일만 업로드할 수 있습니다."
            ),
        )

    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="빈 파일입니다.")

    logger.info(
        "[RagDocumentRouter] received %s (%s, %d bytes)",
        filename,
        file.content_type,
        len(content),
    )

    text = _extract_text(filename, content)
    return await use_case.upload_document(
        filename=filename,
        content_type=file.content_type or "",
        text=text,
    )
